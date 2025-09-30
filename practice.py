#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
PyQt6 Quiz App 
Capabilities:
- Multiple Choice + Multi-Select questions
- Images pulled from PPTX; click image thumbnail to show enlarged image
- Flag questions (toggle flag, jump to flagged list)
- One 15-minute break (enabled only if timer > 0)
- Calculator integrated
- Check current answer (disabled in Test Mode)
- Pick number of questions: all or specific (supports repeats > max)
- Timer (0 to disable timing and break button)
- Allow repeat questions
- Show Reason/Explanation (disabled in Test Mode)
- Build question bank from PowerPoint slides (notes-driven answers)
- Test Mode: disables Check/Reason during the run
- Submit button indicates answer saved
- Check button will indicate correct or not quite if incorrect. 
- Built in dependency checking for required libraries and this program. 
  It not latest version it will prompt you to download the latest version.

| Shorthand  | Longhand meaning                                                                 |
| ---------- | -------------------------------------------------------------------------------- |
| `q`        | The current **question dict** (or question text in review/export loops).         |
| `opts`     | List of **answer option strings** for a question.                                |
| `cb`       | A **QCheckBox** representing one multi-select option.                            |
| `btn`      | Generic **button** object (QPushButton/QRadioButton) used in helpers and rows.   |
| `dlg`      | A **dialog** instance (e.g., `BreakDialog`, flag list `QDialog`).                |
| `v`        | A **QVBoxLayout** (vertical layout) placeholder var.                             |
| `row`      | A **QHBoxLayout** (or row container) placeholder var for footer/action rows.     |
| `hl`       | A **horizontal layout** used inside an answer row.                               |
| `lst`      | The **QListWidget** that lists flagged questions.                                |
| `it`       | The **currently selected item** (`QListWidgetItem`) in the flag dialog.          |
| `idx`      | The **current question index** (int).                                            |
| `mm`, `ss` | **Minutes** / **seconds** when formatting timers.                                |
| `img`      | The **image path** for the current question.                                     |
| `pix`      | A **QPixmap** loaded from `img`.                                                 |
| `scaled`   | The **scaled pixmap** used for the thumbnail label.                              |
| `ms`       | **Milliseconds** for flash/revert timers on buttons.                             |
| `t`        | The in-flight **QTimer** used to revert a flashed button.                        |
| `n`        | **Requested question count** from Settings.                                      |
| `r`        | One **review/export row dict** while writing the .txt.                           |
| `i`        | **Loop index** (1-based in export).                                              |
| `w`        | Generic **widget** variable in “add this widget” loops.                          |
| `cw`       | The **central widget** (`QWidget`) for the main window.                          |
| `root`     | The main window’s **root layout** (`QVBoxLayout`).                               |
| `head`     | The **header row** layout (timer, theme, mode, flags).                           |
| `actions`  | The **actions row** (Show Image, Why, Check, Break…).                            |
| `nav`      | The **navigation row** (Prev/Next/Submit/Finish).                                |
| `_e`       | Throwaway **event object** in mouse handlers.                                    |

"""

from __future__ import annotations

import os
import re
import sys
import random
import subprocess
from dataclasses import dataclass
from typing import List, Dict, Set, Tuple, Optional, Callable

import json
from urllib.request import Request, urlopen

VERSION = "1.2.4"

#---------------------------------
#  Checking for required libraries
#---------------------------------

try:
    from importlib.metadata import version, PackageNotFoundError  # Python 3.8+
except Exception:  # pragma: no cover (older Pythons)
    from importlib_metadata import version, PackageNotFoundError  # type: ignore


def _parse_version(v: str) -> tuple:
    """Lightweight version tuple (major, minor, patch) for comparisons."""
    parts = v.split(".")
    nums = []
    for p in parts[:3]:
        try:
            nums.append(int("".join(ch for ch in p if ch.isdigit())))
        except ValueError:
            nums.append(0)
    while len(nums) < 3:
        nums.append(0)
    return tuple(nums)


def _get_installed_version(dist_name: str) -> Optional[str]:
    """Return installed distribution version (or None if not installed)."""
    try:
        return version(dist_name)
    except PackageNotFoundError:
        return None


def _needs(dist_name: str, min_ver_str: str) -> bool:
    """Return True if dist_name is missing or below min_ver_str."""
    v = _get_installed_version(dist_name)
    if v is None:
        return True
    return _parse_version(v) < _parse_version(min_ver_str)


def _ask_yes_no(prompt: str, default_no: bool = True) -> bool:
    """Return True if user answers yes. If not interactive, return False when default_no."""
    if not sys.stdin or not sys.stdin.isatty():
        return not default_no
    try:
        ans = input(prompt).strip().lower()
    except EOFError:
        return not default_no
    return ans in ("y", "yes")


def _pip_install(args: List[str]) -> None:
    """Run pip install/upgrade with provided args (list of 'pkg>=x.y')."""
    cmd = [sys.executable, "-m", "pip", "install", "--upgrade", *args]
    subprocess.check_call(cmd)


def _force_popup_update_warning(download_page_url: str, parent=None) -> None:
    """
    Always show a QMessageBox warning. If no QApplication exists yet,
    create a temporary one just for this modal dialog.
    """
    try:
        from PyQt6.QtWidgets import QApplication, QMessageBox
        from PyQt6.QtGui import QDesktopServices
        from PyQt6.QtCore import QUrl
    except Exception:
        # If PyQt isn't available for some reason, silently return.
        return

    app_created = False
    app = QApplication.instance()
    if app is None:
        app = QApplication(sys.argv)
        app_created = True

    box = QMessageBox(parent)
    box.setIcon(QMessageBox.Icon.Warning)
    box.setWindowTitle("Update Available")
    box.setText(
        "Your version is not the latest version.\n\n"
        "For the latest version and features please download a new version from\n"
        f"{download_page_url}"
    )
    open_button = box.addButton("Open Download Page", QMessageBox.ButtonRole.AcceptRole)
    box.addButton("OK", QMessageBox.ButtonRole.RejectRole)
    box.exec()

    if box.clickedButton() is open_button:
        QDesktopServices.openUrl(QUrl(download_page_url))


def ensure_requirements() -> None:
    """
    Ensure required packages (PyQt6, python-pptx) are present at minimum versions.
    Prompts once (if interactive). Installs/updates both in one pip call when needed.
    Exits with a helpful message if requirements remain unsatisfied.
    Then compare local version of program with github.
    """
    requirements = [
        ("PyQt6",       "6.5.0"),   # adjust if you need a newer baseline
        ("python-pptx", "0.6.21"),
    ]
    
    to_install = [f"{name}>={minver}" for (name, minver) in requirements if _needs(name, minver)]

    if not to_install:
        # Optional: print once for verbose mode
        #print("✅ All dependencies satisfied.")
        pass
        
    else:
        # One prompt for all missing/outdated packages
        message_lines = [
            "The following packages are missing or below the required version:",
            *[f"  - {spec}" for spec in to_install],
            "",
        ]
        print("\n".join(message_lines))

        if _ask_yes_no("Install/upgrade them now? (y/n): ", default_no=True):
            try:
                print("Installing/upgrading: " + ", ".join(to_install))
                _pip_install(to_install)
            except subprocess.CalledProcessError as e:
                print(f"❌ Installation failed: {e}")
                sys.exit(1)
        else:
            print("❌ Required packages not installed. You can install them with:")
            print(f"   {sys.executable} -m pip install --upgrade " + " ".join(to_install))
            sys.exit(1)

        # Re-verify after pip run
        still_missing = [f"{n}>={v}" for (n, v) in requirements if _needs(n, v)]
        if still_missing:
            print("❌ Some requirements are still unsatisfied after installation:")
            for spec in still_missing:
                print("  -", spec)
            sys.exit(1)
    
    # -----------------------------
    # GitHub VERSION file check
    # -----------------------------
    
    # Config: your raw VERSION file and the page you want to open
    GITHUB_VERSION_FILE_URL = "https://raw.githubusercontent.com/Comrob2018/QuizApp/main/VERSION"
    GITHUB_DOWNLOAD_PAGE    = "https://github.com/Comrob2018/QuizApp/tree/main"

    try:

        def _fetch_remote_text(url: str, timeout: float = 4.0) -> str | None:
            try:
                req = Request(url, headers={
                    "Accept": "text/plain, */*;q=0.1",
                    "User-Agent": "quiz-app-version-check"
                })
                with urlopen(req, timeout=timeout) as resp:
                    return resp.read().decode("utf-8", errors="replace")
            except Exception:
                return None

        def _extract_remote_version(text: str) -> str | None:
            if not text:
                return None
            first = text.strip().splitlines()[0].strip()
            m = re.search(r"\b\d+(?:\.\d+){1,3}[A-Za-z0-9\-\.]*\b", first)
            if m:
                return m.group(0)
            # JSON fallback
            try:
                obj = json.loads(text)
                cand = obj.get("latest") or obj.get("version") or obj.get("tag_name")
                if isinstance(cand, str):
                    return cand[1:] if cand.lower().startswith("v") else cand
            except Exception:
                pass
            return None

        remote_text = _fetch_remote_text(GITHUB_VERSION_FILE_URL)
        remote_ver  = _extract_remote_version(remote_text) if remote_text else None

        if remote_ver and _parse_version(VERSION) < _parse_version(remote_ver):
            _force_popup_update_warning(GITHUB_DOWNLOAD_PAGE)

    except Exception:
        # Swallow any network/parse errors to avoid blocking startup.
        pass

from PyQt6.QtWidgets import (
    QApplication, QMainWindow, QWidget, QLabel, QPushButton, QRadioButton,
    QVBoxLayout, QHBoxLayout, QGridLayout, QDialog, QMessageBox, QFileDialog,
    QScrollArea, QButtonGroup, QCheckBox, QLineEdit, QFrame, QSizePolicy,
    QListWidget, QListWidgetItem, QComboBox,
)
from PyQt6.QtCore import Qt, QTimer, QSize, QSettings
from PyQt6.QtGui import QPixmap, QFont, QFontDatabase

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

# -----------------------------
# Data model
# -----------------------------

@dataclass
class QuizItem:
    question: str
    options: List[str]
    answer: Set[str]                 # set of correct option strings
    explanation: str
    image: Optional[str] = None
    multi: bool = False


def normalize_quiz_records(raw_items: List[Dict]) -> List[Dict]:
    normalized: List[Dict] = []
    for item in raw_items or []:
        q = str(item.get("question", "")).strip()
        opts = item.get("options") or item.get("answers") or []
        if isinstance(opts, dict):
            opts = list(opts.values())
        opts = [str(o).strip() for o in opts if str(o).strip()]

        correct_raw = item.get("answer") or set()
        if isinstance(correct_raw, str):
            correct = {correct_raw.strip()} if correct_raw.strip() else set()
        elif isinstance(correct_raw, (list, set, tuple)):
            correct = {str(x).strip() for x in correct_raw if str(x).strip()}
        else:
            correct = set()

        explanation = str(item.get("explanation", "")).strip()
        image_path = item.get("image") or None
        image_path = str(image_path) if image_path else None
        multi = bool(item.get("multi", len(correct) > 1))

        normalized.append({
            "question": q,
            "options": opts,
            "answer": correct,
            "explanation": explanation,
            "image": image_path,
            "multi": multi
        })
    return normalized


# -----------------------------
# PPTX parser (notes-driven answers)
# -----------------------------

_CORRECT_SEP_RE = re.compile(r"\s*[;|]\s*")  # split on ';' or '|'

def _clean_option_line(line: str) -> str:
    s = (line or "").strip()
    s = re.sub(r"^[\u2022\-–•]\s+", "", s)             # bullets
    s = re.sub(r"^([A-Za-z][\)\.]|\d+\.)\s+", "", s)   # A) A. 1. etc.
    return s.strip()

def _text_frame_lines(text: str) -> List[str]:
    return [ln.strip() for ln in (text or "").splitlines() if ln.strip()]

def _collect_question_and_options_from_shapes(slide) -> Tuple[str, List[str]]:
    """Question: first text shape; Options: subsequent text shapes/lines."""
    question = ""
    options: List[str] = []
    for shp in slide.shapes:
        if getattr(shp, "has_text_frame", False):
            txt = (shp.text or "").strip()
            if not txt:
                continue
            if not question:
                question = txt
                continue
            lines = _text_frame_lines(txt)
            if len(lines) <= 1:
                opt = _clean_option_line(txt)
                if opt:
                    options.append(opt)
            else:
                for ln in lines:
                    opt = _clean_option_line(ln)
                    if opt:
                        options.append(opt)
    return question, options

def _read_notes_answer_and_reason(slide) -> Tuple[Set[str], str]:
    """
    Notes expected:
      Answer is: token1 | token2 ; token3
      Reason on next non-empty line
    """
    try:
        notes = slide.notes_slide.notes_text_frame.text or ""
    except Exception:
        notes = ""
    lines = [ln.strip() for ln in notes.splitlines()]
    answer_line_idx = -1
    answer_str = ""
    for idx, ln in enumerate(lines):
        m = re.search(r"answer\s*is\s*:\s*(.+)", ln, flags=re.IGNORECASE)
        if m:
            answer_line_idx = idx
            answer_str = m.group(1).strip()
            break
    reason = ""
    if answer_line_idx != -1:
        for j in range(answer_line_idx + 1, len(lines)):
            if lines[j]:
                reason = lines[j]
                break
    tokens = [t.strip() for t in _CORRECT_SEP_RE.split(answer_str) if t.strip()]
    return set(tokens), reason


def _map_correct_tokens_to_options(tokens: Set[str], options: List[str]) -> Set[str]:
    """Exact-match-first; fallback to substring; supports letter mapping."""
    if not tokens:
        return set()
    mapped: Set[str] = set()
    letter_map = {chr(ord('A') + i): opt for i, opt in enumerate(options)}
    low_options = [(opt, opt.lower()) for opt in options]
    for t in tokens:
        t_str = str(t).strip()
        if not t_str:
            continue
        key = t_str.upper()
        if key in letter_map:
            mapped.add(letter_map[key])
            continue
        tl = t_str.lower()
        exact = next((opt for opt, low in low_options if low == tl), None)
        if exact:
            mapped.add(exact)
            continue
        sub = next((opt for opt, low in low_options if tl in low or low in tl), None)
        if sub:
            mapped.add(sub)
    return mapped


def extract_images_and_prepare_quiz(pptx_path: str, out_dir: Optional[str] = None) -> List[Dict]:
    prs = Presentation(pptx_path)
    stem = os.path.splitext(os.path.basename(pptx_path))[0]
    base_out = out_dir or os.path.join(os.path.dirname(pptx_path), f"extracted_images")
    os.makedirs(base_out, exist_ok=True)

    quiz_items: List[Dict] = []

    for s_idx, slide in enumerate(prs.slides, start=1):
        question_text, options = _collect_question_and_options_from_shapes(slide)
        if not question_text or not options:
            continue

        raw_tokens, reason = _read_notes_answer_and_reason(slide)
        correct_set = _map_correct_tokens_to_options(raw_tokens, options)

        # First picture on slide becomes question image
        q_image_path = None
        image_counter = 0
        for shp in slide.shapes:
            if shp.shape_type == MSO_SHAPE_TYPE.PICTURE:
                image_counter += 1
                image = shp.image
                ext = image.ext or "png"
                image_name = f"{stem}-slide-{s_idx:02d}-{image_counter:02d}.{ext}"
                out_path = os.path.join(base_out, image_name)
                with open(out_path, "wb") as f:
                    f.write(image.blob)
                if not q_image_path:
                    q_image_path = out_path

        item = {
            "question": question_text.strip(),
            "options": options,
            "answer": correct_set,          # set[str]
            "explanation": reason.strip(),
            "image": q_image_path,
            # multi if notes had multiple tokens OR mapped to multiple options
            "multi": (len(raw_tokens) > 1) or (len(correct_set) > 1),
        }
        quiz_items.append(item)

    return normalize_quiz_records(quiz_items)


def build_quiz_from_pptx(pptx_path: str) -> Tuple[List[Dict], str]:
    data = extract_images_and_prepare_quiz(pptx_path)
    return data, pptx_path

# -----------------------------
#       --- THEME QSS ---
# -----------------------------

THEMES = {
    "dark": {
        "bg": "#1c1c1c", "surface": "#2a2a2a", "surface_alt": "#333333", "text": "#eaeaea",
        "muted": "#777777", "border": "#2f2f2f", "primary": "#7AA2F7", "accent": "#BB9AF7",
        "success": "#9ECE6A", "warn": "#E0AF68", "error": "#F7768E"
    },
    "solarized_dark": {
        "bg": "#002b36", "surface": "#073642", "surface_alt": "#0a3a46", "text": "#eee8d5",
        "muted": "#93a1a1", "border": "#586e75", "primary": "#268bd2", "accent": "#6c71c4",
        "success": "#859900", "warn": "#b58900", "error": "#dc322f"
    },
    "nord": {
        "bg": "#2E3440", "surface": "#3B4252", "surface_alt": "#434C5E", "text": "#ECEFF4",
        "muted": "#D8DEE9", "border": "#434C5E", "primary": "#88C0D0", "accent": "#5E81AC",
        "success": "#A3BE8C", "warn": "#EBCB8B", "error": "#BF616A"
    },
    "gruvbox_dark": {
        "bg": "#282828", "surface": "#3C3836", "surface_alt": "#504945", "text": "#EBDBB2",
        "muted": "#A89984", "border": "#A89984", "primary": "#83A598", "accent": "#D3869B",
        "success": "#B8BB26", "warn": "#FABD2F", "error": "#FB4934"
    },
    "tokyo_night": {
        "bg": "#1A1B26", "surface": "#292E42", "surface_alt": "#2f3353", "text": "#C0CAF5",
        "muted": "#9AA5CE", "border": "#3B4261", "primary": "#7AA2F7", "accent": "#BB9AF7",
        "success": "#9ECE6A", "warn": "#E0AF68", "error": "#F7768E"
    },
    "high_contrast": {
        "bg": "#000000", "surface": "#111111", "surface_alt": "#1A1A1A", "text": "#FFFFFF",
        "muted": "#BFBFBF", "border": "#FFFFFF", "primary": "#FFD400", "accent": "#00B8FF",
        "success": "#00E5A0", "warn": "#FFB000", "error": "#FF3B30"
    },
    "cyberpunk": {
        "bg": "#1A002B", "surface": "#24033A", "surface_alt": "#2F0B4A", "text": "#FCEFFF",
        "muted": "#CAAEDF", "border": "#4A2A6A", "primary": "#00F0FF", "accent": "#FF6BD6",
        "success": "#89FFBF", "warn": "#FFC857", "error": "#FF5C8A"
    },
    "sapphire": {
        "bg": "#162C45", "surface": "#353652", "surface_alt": "#4C334D", "text": "#F2F2F2",  
        "muted": "#A9A4AF", "border": "#4C334D", "primary": "#053C5E", "accent": "#DB222A",  
        "success": "#4CAF50", "warn": "#FF9800", "error": "#C32530",  
    },
    "crimson_ember": {
        "bg": "#03071E", "surface": "#370617", "surface_alt": "#6A040F", "text": "#FFF6E5",  
        "muted": "#CBBDAA", "border": "#9D0208", "primary": "#FFBA08", "accent": "#D00000",  
        "success": "#4CAF50", "warn": "#F48C06", "error": "#D00000",  
    },
     "aurora_teal": {
        "bg": "#0B1020", "surface": "#121A2B", "surface_alt": "#1A2338", "text": "#E8F6FF",
        "muted": "#9FB3C8", "border": "#253352", "primary": "#2DD4BF", "accent": "#7C3AED",  
        "success": "#22C55E", "warn": "#F59E0B", "error": "#EF4444",
    },
    "dark_rose": {
        "bg": "#450A1B", "surface": "#800F2F", "surface_alt": "#A4133C", "text": "#FFF0F3", 
        "muted": "#FF8FA3", "border": "#A4133C", "primary": "#FF758F", "accent": "#FF4D6D",
        "success": "#4CAF50", "warn": "#FF9800", "error": "#C9184A",
    },
    "hacker_terminal": {
        "bg": "#0A0F0A","surface": "#0F1A0F","surface_alt": "#152415","text": "#E8FFE8",
        "muted": "#89A989","border": "#1E2F1E","primary": "#00FF7A","accent": "#00F0FF",
        "success": "#00E676","warn": "#FFC857","error": "#FF4D6D",
    },
}

# Human-friendly names -> keys (what shows in the dropdown)
THEME_NAMES = {
    "Dark": "dark",
    "Solarized Dark": "solarized_dark",
    "Nord": "nord",
    "Gruvbox (Dark)": "gruvbox_dark",
    "Tokyo Night": "tokyo_night",
    "High Contrast": "high_contrast",
    "Cyberpunk": "cyberpunk",
    "Sapphire": "sapphire",
    "Crimson Ember": "crimson_ember",
    "Aurora": "aurora_teal",
    "Dark Rose": "dark_rose",
    "Hacker Terminal": "hacker_terminal",
}

# Convenience: inverse map if you ever need to go from key -> label
THEME_LABELS = {v: k for k, v in THEME_NAMES.items()}

_TERMINAL_FONT_CANDIDATES = [
    # nice modern monos first
    "Cascadia Mono", "JetBrains Mono", "Fira Code", "IBM Plex Mono", "Hack",
    # common system monos
    "Consolas", "Menlo", "Monaco", "Lucida Console",
    "Liberation Mono", "DejaVu Sans Mono", "Source Code Pro",
]

_DEFAULT_APP_FONT = None  # captured the first time apply_theme runs

def _pick_first_available_font(candidates=_TERMINAL_FONT_CANDIDATES) -> str:
    try:
        # PyQt6: families() is a class method; no instance needed
        installed = set(QFontDatabase.families())
        for name in candidates:
            if name in installed:
                return name
        # Reliable fallback to system fixed-width font
        return QFontDatabase.systemFont(QFontDatabase.SystemFont.FixedFont).family()
    except Exception:
        # Last-ditch fallback if called too early (before QApplication) or in odd envs
        return "Monospace"


def _ideal_on(hex_color: str) -> str:
    c = hex_color.lstrip("#")
    r, g, b = int(c[0:2],16), int(c[2:4],16), int(c[4:6],16)
    y = (0.2126*r + 0.7152*g + 0.0722*b)
    return "#000000" if y > 186 else "#FFFFFF"

def _build_stylesheet_from_theme(t: dict) -> str:
    return f"""
    QWidget {{
        background-color: {t['bg']};
        color: {t['text']};
    }}
    QLabel#Heading {{ font-size: 16px; font-weight: 600; }}
    QFrame[card="true"], QFrame#Card, QLineEdit, QTextEdit, QPlainTextEdit, QListWidget {{
        background-color: {t['surface']};
        border: 1px solid {t['border']};
        border-radius: 8px;
    }}
    QPushButton {{
        background-color: {t['surface']};
        border: 1px solid {t['border']};
        border-radius: 10px;
        padding: 6px 12px;
        color: {t['text']};
    }}
    QPushButton:hover {{ background-color: {t['surface_alt']}; }}
    QPushButton:disabled {{ color: {t['muted']}; }}
    QPushButton#primary {{
        background-color: {t['primary']};
        color: {_ideal_on(t['primary'])};
        border: none;
    }}
    QPushButton#danger {{
        background-color: {t['error']};
        color: {_ideal_on(t['error'])};
        border: none;
    }}
    QPushButton:focus, QLineEdit:focus, QComboBox:focus, QListWidget:focus, QTextEdit:focus, QPlainTextEdit:focus {{
        border: 2px solid {t['primary']};
    }}
    QListWidget::item:selected {{
        background: {t['primary']};
        color: {_ideal_on(t['primary'])};
    }}
    QCheckBox::indicator {{
        width: 18px;
        height: 18px;
        border-radius: 4px;
        border: 1px solid {t['border']};
        background: {t['surface']};
    }}
    QCheckBox::indicator:hover {{
        border-color: {t['primary']};
    }}
    QCheckBox::indicator:checked {{
        background: {t['primary']};
        border-color: {t['primary']};
    }}
    QCheckBox::indicator:indeterminate {{
        background: {t['accent']};
        border-color: {t['accent']};
    }}
    QRadioButton::indicator {{
        width: 18px; height: 18px;
        border-radius: 9px;
        border: 1px solid {t['border']};
        background: {t['surface']};
    }}
    QRadioButton::indicator:checked {{
        background: {t['primary']};
        border-color: {t['primary']};
    }}
    QLabel#MultiHint{{
        font-size: 12px;
        color: {t['muted']};
        margin: 2px 0 6px 0;
    }}
    QLabel#ProgressLabel {{
        font-size: 13px;
        font-weight: 600;
        color: {t['muted']};
    }}
    """

def apply_theme(app: "QApplication", theme_name: str):
    global _DEFAULT_APP_FONT
    if _DEFAULT_APP_FONT is None:
        _DEFAULT_APP_FONT = app.font()

    # fallback if unknown
    if theme_name not in THEMES:
        theme_name = "tokyo_night"
    t = THEMES[theme_name]

    # Palette for native widgets (selection colors, text, etc.)
    from PyQt6.QtGui import QPalette, QColor
    pal = QPalette()
    pal.setColor(QPalette.ColorRole.Window, QColor(t["bg"]))
    pal.setColor(QPalette.ColorRole.Base, QColor(t["surface"]))
    pal.setColor(QPalette.ColorRole.AlternateBase, QColor(t["surface_alt"]))
    pal.setColor(QPalette.ColorRole.WindowText, QColor(t["text"]))
    pal.setColor(QPalette.ColorRole.Text, QColor(t["text"]))
    pal.setColor(QPalette.ColorRole.Button, QColor(t["surface"]))
    pal.setColor(QPalette.ColorRole.ButtonText, QColor(t["text"]))
    pal.setColor(QPalette.ColorRole.Highlight, QColor(t["primary"]))
    pal.setColor(QPalette.ColorRole.HighlightedText, QColor(_ideal_on(t["primary"])))
    app.setPalette(pal)

    # Stylesheet: build from tokens, then (for your legacy light/dark) append your QSS.
    base_qss = _build_stylesheet_from_theme(t)
    # Hacker Terminal → switch to a classic monospace; restore for others
    if theme_name == "hacker_terminal" or theme_name == "cyberpunk":
        fam = _pick_first_available_font()
        app.setFont(QFont(fam, 10))                   
        base_qss += f'* {{ font-family: "{fam}"; }}\n' # enforce monospace everywhere
    else:
        if _DEFAULT_APP_FONT is not None:
            app.setFont(_DEFAULT_APP_FONT)


    base_qss += f"""
    QPushButton#primary {{ color: {_ideal_on(t['primary'])}; }}
    QPushButton#danger  {{ color: {_ideal_on(t['error'])};  }}
    """
    app.setStyleSheet(base_qss)
    QSettings("YourOrg", "QuizApp").setValue("theme", theme_name)

def load_theme_pref() -> str:
    v = QSettings("YourOrg", "QuizApp").value("theme", "dark")
    return v if isinstance(v, str) and v in THEMES else "dark"

# -----------------------------
# Dialogs: Image viewer, Calculator, Review, Settings
# -----------------------------

class ImageViewerDialog(QDialog):
    def __init__(self, image_paths: List[str], parent: Optional[QWidget] = None, title: str = "Images"):
        super().__init__(parent)
        self.setWindowTitle(title)
        self.resize(800, 600)
        outer = QVBoxLayout(self)
        scroll = QScrollArea(self)
        scroll.setWidgetResizable(True)
        outer.addWidget(scroll)
        container = QWidget()
        v = QVBoxLayout(container)
        v.setContentsMargins(0,0,0,0)
        v.setSpacing(8)
        for path in image_paths:
            if not os.path.exists(path):
                label = QLabel(f"(Missing image) {path}")
                label.setStyleSheet("color:#b00;font-style:italic;")
                v.addWidget(label)
                continue
            pix = QPixmap(path)
            image_label = QLabel()
            image_label.setPixmap(pix)
            image_label.setScaledContents(True)
            image_label.setMinimumSize(QSize(400, min(500, pix.height())))
            v.addWidget(image_label)
            cap = QLabel(os.path.basename(path))
            cap.setAlignment(Qt.AlignmentFlag.AlignLeft)
            cap.setStyleSheet("color:#555;")
            v.addWidget(cap)
            line = QFrame()
            line.setFrameShape(QFrame.Shape.HLine)
            v.addWidget(line)
        v.addStretch(1)
        scroll.setWidget(container)
        row = QHBoxLayout()
        row.addStretch(1)
        close_button = QPushButton("Close")
        close_button.clicked.connect(self.accept)
        row.addWidget(close_button)
        outer.addLayout(row)


class CalculatorPopup(QDialog):
    def __init__(self, parent: Optional[QWidget] = None):
        super().__init__(parent)
        self.setWindowTitle("Calculator")
        self.expr = ""
        outer = QVBoxLayout(self)
        self.display = QLineEdit()
        self.display.setReadOnly(True)
        outer.addWidget(self.display)
        grid = QGridLayout()
        buttons = [
            ("7",0,0),("8",0,1),("9",0,2),("/",0,3),
            ("4",1,0),("5",1,1),("6",1,2),("*",1,3),
            ("1",2,0),("2",2,1),("3",2,2),("-",2,3),
            ("0",3,0),(".",3,1),("C",3,2),("+",3,3),
        ]
        for text, r, c in buttons:
            b = QPushButton(text)
            b.clicked.connect(lambda _, t=text: self.click_button(t))
            grid.addWidget(b, r, c)
        outer.addLayout(grid)
        row = QHBoxLayout()
        eq = QPushButton("=")
        eq.clicked.connect(self.evaluate)
        clr = QPushButton("Clear")
        clr.clicked.connect(self.clear_display)
        row.addWidget(eq)
        row.addStretch(1)
        row.addWidget(clr)
        outer.addLayout(row)

    def click_button(self, t: str):
        if t == "C":
            self.clear_display()
        else:
            self.expr += t
            self.display.setText(self.expr)

    def clear_display(self):
        self.expr = ""
        self.display.clear()

    def evaluate(self):
        try:
            result = str(eval(self.expr, {"__builtins__": {}}, {}))
            self.expr = result
            self.display.setText(result)
        except Exception:
            QMessageBox.warning(self, "Error", "Invalid expression.")


class ReviewPopup(QDialog):
    """Review with percentage, built when Finish is clicked.

    Args:
        parent: Parent widget (use None to make it a top-level window).
        review_items: List of dicts with keys:
            - "question": str
            - "correct": str
            - "chosen": str
            - "explanation": str
            - "flagged": bool
        score_tuple: (correct_count, total_count, percent_int)
        pptx_basename: Used to name exported files.
        restart_callback: Optional[Callable[[], None]] to start a new test run.
    """
    def __init__(
        self,
        parent: Optional[QWidget],
        review_items: List[Dict],
        score_tuple: Tuple[int, int, int],
        pptx_basename: str = "review",
        restart_callback: Optional[Callable[[], None]] = None
    ):
        super().__init__(parent)
        self.review_items = review_items
        self.correct, self.total, self.percent = score_tuple
        self.pptx_basename = os.path.splitext(os.path.basename(pptx_basename))[0] if pptx_basename else "review"
        self._restart_callback = restart_callback

        self.setWindowTitle("Review")
        self.resize(900, 640)

        outer = QVBoxLayout(self)

        # Score header
        hdr = QLabel(f"Score: {self.correct}/{self.total} ({self.percent}%)")
        hdr.setStyleSheet("font-weight:600; font-size:14px;")
        outer.addWidget(hdr)

        # Scrollable body
        scroll = QScrollArea(self)
        scroll.setWidgetResizable(True)
        outer.addWidget(scroll)

        container = QWidget()
        v = QVBoxLayout(container)
        v.setContentsMargins(8, 8, 8, 8)
        v.setSpacing(8)

        for i, r in enumerate(self.review_items, start=1):
            q = r.get("question", "")
            correct = r.get("correct", "")
            chosen = r.get("chosen", "")
            reason = r.get("explanation", "")

            is_incorrect = chosen.strip() and (chosen != correct)

            text = (
                f"{i}. {q}\n"
                f"correct answer: {correct}\n"
                f"your answer: {chosen}\n"
                f"explanation: {reason}\n"
            )

            label = QLabel(text)
            label.setWordWrap(True)

            if is_incorrect:
                # Add a red ❌ marker and background highlight
                label.setText(f"❌ {text}")
            else:
                # Mark correct answers with a subtle ✅
                label.setText(f"✅ {text}")

            v.addWidget(label)

            line = QFrame()
            line.setFrameShape(QFrame.Shape.HLine)
            line.setFrameShadow(QFrame.Shadow.Sunken)
            v.addWidget(line)

        v.addStretch(1)
        scroll.setWidget(container)

        # Footer row
        row = QHBoxLayout()
        export_button = QPushButton("Export .txt")
        export_button.clicked.connect(self._export_txt)

        restart_button = QPushButton("Restart")
        def _do_restart():
            if self._restart_callback:
                try:
                    self._restart_callback()
                except Exception as e:
                    QMessageBox.critical(self, "Restart failed", str(e))
            self.accept()
        restart_button.clicked.connect(_do_restart)

        close_button = QPushButton("Close")
        close_button.clicked.connect(self.accept)

        row.addWidget(export_button)
        row.addStretch(1)
        row.addWidget(restart_button)
        row.addWidget(close_button)
        outer.addLayout(row)

    def _export_txt(self):
        """Export the review (including score header and ✓/✗ markers) to a .txt file."""
        from datetime import datetime
        default_name = f"{self.pptx_basename}_{datetime.now().strftime('%Y%m%d_%H%M%S')}.txt"
        path, _ = QFileDialog.getSaveFileName(self, "Save Review", default_name, "Text Files (*.txt)")
        if not path:
            return
        try:
            with open(path, "w", encoding="utf-8") as f:
                f.write(f"Score: {self.correct}/{self.total} ({self.percent}%)\n\n")
                for i, r in enumerate(self.review_items, start=1):
                    q = r.get("question", "")
                    correct = r.get("correct", "")
                    chosen = r.get("chosen", "")
                    reason = r.get("explanation", "")

                    # decide if correct/incorrect
                    is_incorrect = chosen.strip() and (chosen != correct)
                    marker = "✗" if is_incorrect else "✓"

                    f.write(
                        f"{marker} {i}. {q}\n"
                        f"correct answer: {correct}\n"
                        f"your answer: {chosen}\n"
                        f"explanation: {reason}\n\n"
                    )
            QMessageBox.information(self, "Exported", f"Saved to:\n{path}")
        except Exception as e:
            QMessageBox.critical(self, "Error", f"Failed to save file:\n{e}")


class QuestionPopup(QDialog):
    """Settings dialog"""
    def __init__(self, parent: Optional[QWidget], max_questions: int):
        super().__init__(parent)
        self.max_questions = max_questions
        self.result = None
        self.setWindowTitle("Settings")
        self.setModal(True)
        self.resize(520, 260)

        outer = QVBoxLayout(self)
        outer.setSpacing(8)
        outer.setContentsMargins(8,8,8,8)

        row1 = QHBoxLayout()
        row1.addWidget(QLabel(f"""Number of questions (0=all, up to {max_questions}, 
for >{max_questions} you have to allow repeat questions):"""))
        self.count_edit = QLineEdit("0")
        row1.addWidget(self.count_edit)
        outer.addLayout(row1)

        row2 = QHBoxLayout()
        row2.addWidget(QLabel("Timer (mm:ss or mm; 0 = no timer):"))
        self.timer_edit = QLineEdit("0")
        row2.addWidget(self.timer_edit)
        outer.addLayout(row2)

        self.repeat_cb = QCheckBox("Allow repeated questions")
        self.repeat_cb.setChecked(True)
        self.allow_calc_cb = QCheckBox("Allow Calculator")
        self.allow_calc_cb.setChecked(False)
        self.test_mode_cb = QCheckBox("Enable Test Mode (disables Check/Why during test)")
        self.test_mode_cb.setChecked(False)
        self.breaks_cb = QCheckBox("Allow a single 15-minute break (requires timer > 0)")
        self.breaks_cb.setChecked(True)
        for w in (self.repeat_cb, self.allow_calc_cb, self.test_mode_cb, self.breaks_cb):
            outer.addWidget(w)

        buttons = QHBoxLayout()
        buttons.addStretch(1)
        ok = QPushButton("Start")
        cancel = QPushButton("Cancel")
        ok.clicked.connect(self.on_ok)
        cancel.clicked.connect(self.reject)
        buttons.addWidget(ok); buttons.addWidget(cancel)
        outer.addLayout(buttons)

    def _parse_timer(self, text: str) -> int:
        text = (text or "").strip()
        if not text or text == "0":
            return 0
        if ":" in text:
            mm, ss = text.split(":", 1)
            try:
                return int(mm) * 60 + int(ss)
            except ValueError:
                return 0
        try:
            return int(text) * 60
        except ValueError:
            return 0

    def on_ok(self):
        raw = self.count_edit.text().strip() or "0"
        try:
            n = int(raw)
        except ValueError:
            QMessageBox.warning(self, "Invalid number", "Enter a whole number for question count.")
            return
        if n < 0:
            QMessageBox.warning(self, "Invalid number", "Question count cannot be negative.")
            return
        timer_seconds = self._parse_timer(self.timer_edit.text())
        self.result = (
            n,
            timer_seconds,
            self.repeat_cb.isChecked(),
            self.allow_calc_cb.isChecked(),
            self.test_mode_cb.isChecked(),
            self.breaks_cb.isChecked()
        )
        self.accept()

    def get_result(self):
        return self.result


class BreakDialog(QDialog):
    """15-minute break dialog with live countdown."""
    def __init__(self, parent=None, total_seconds: int = 15 * 60):
        super().__init__(parent)
        self.setWindowTitle("Break")
        self.setModal(True)
        self.resize(340, 150)

        self.remaining = int(total_seconds)
        self.timer = QTimer(self)
        self.timer.setInterval(1000)
        self.timer.timeout.connect(self._tick)

        v = QVBoxLayout(self)
        self.title = QLabel("Break in progress")
        self.title.setAlignment(Qt.AlignmentFlag.AlignCenter)
        self.title.setStyleSheet("font-weight: 600; font-size: 14px;")
        v.addWidget(self.title)

        self.time_label = QLabel("")
        self.time_label.setAlignment(Qt.AlignmentFlag.AlignCenter)
        self.time_label.setStyleSheet("font-family: monospace; font-size: 18px;")
        v.addWidget(self.time_label)

        row = QHBoxLayout()
        row.addStretch(1)
        self.end_button = QPushButton("End Break Now")
        self.end_button.clicked.connect(self.accept)  # closes dialog
        row.addWidget(self.end_button)
        v.addLayout(row)

        self._update_label()
        self.timer.start()

    def _tick(self):
        self.remaining -= 1
        if self.remaining <= 0:
            self.timer.stop()
            self.remaining = 0
            self._update_label()
            self.accept()  # auto-close when timer reaches 0
        else:
            self._update_label()

    def _update_label(self):
        mm, ss = divmod(max(0, self.remaining), 60)
        self.time_label.setText(f"{mm:02d}:{ss:02d}")

# -----------------------------
# Main Window
# -----------------------------

class QuizMainWindow(QMainWindow):
    def __init__(self, quiz_data: List[Dict], pptx_path: Optional[str] = None, settings: Optional[Tuple]=None):
        super().__init__()
        self.raw_quiz: List[Dict] = quiz_data[:]
        self.pptx_path = pptx_path or "review"
        self.flags: Set[int] = set()      # flagged question indices (in current order)
        self.break_taken = False
        self.test_mode = False
        self.allow_breaks = True
        self.allow_calc = False

        # Apply settings
        count = 0; timer_seconds = 0; allow_repeats = True; test_mode = False; allow_breaks = True
        if settings:
            count, timer_seconds, allow_repeats, allow_calc, test_mode, allow_breaks = settings
        self.test_mode = bool(test_mode)
        self.allow_breaks = bool(allow_breaks)
        self.allow_calc = bool(allow_calc)

        self.quiz = self._apply_quiz_settings(self.raw_quiz, count, allow_repeats)
        # Randomize the order of answers ONCE at quiz start
        for q in self.quiz:
            opts = q.get("options", [])
            if opts:
                random.shuffle(opts)
        self.user_answers: List[Optional[Set[str] | str]] = [None] * len(self.quiz)
        self.setWindowTitle("Quiz App")
        self.resize(980, 700)

        # Timer
        self.total_seconds = max(0, int(timer_seconds))
        self.remaining_seconds = self.total_seconds
        self.timer = QTimer(self)
        self.timer.timeout.connect(self._on_tick)

        # --- UI layout ---
        cw = QWidget(); self.setCentralWidget(cw)
        root = QVBoxLayout(cw); root.setSpacing(8); root.setContentsMargins(8,8,8,8)

        # Header
        head = QHBoxLayout(); head.setSpacing(8)
        self.timer_label = QLabel("")
        self.timer_label.setStyleSheet("font-weight: 600; font-size: 18px; letter-spacing: 1px;")
        self.progress_label = QLabel("0/0")
        self.progress_label.setStyleSheet("font-weight: 600; font-size: 18px; letter-spacing: 1px;")
        self.flag_button = QPushButton("Flag")
        self.flag_button.setToolTip("Toggle flag for this question")
        self.flag_button.clicked.connect(self._toggle_flag)
        self.flag_list_button = QPushButton("Flagged…")
        
        # Theme picker (dropdown)
        self.theme_combo = QComboBox()
        self.theme_combo.setEditable(False)
        self.theme_combo.setInsertPolicy(QComboBox.InsertPolicy.NoInsert)

        # Populate with (label, key)
        for label, key in sorted(THEME_NAMES.items(), key=lambda x: x[0].lower()):
            self.theme_combo.addItem(label, userData=key)

        current_key = load_theme_pref()
        # Set visible label based on stored key
        idx = self.theme_combo.findData(current_key)
        if idx >= 0:
            self.theme_combo.setCurrentIndex(idx)

        self.theme_combo.setToolTip("Choose a theme")
        self.theme_combo.currentIndexChanged.connect(self._on_theme_changed)

        # --- Mode badge (Study vs Test) ---
        self.mode_badge = QLabel()
        self.mode_badge.setAlignment(Qt.AlignmentFlag.AlignCenter)
        self.mode_badge.setStyleSheet(
            "QLabel { padding: 3px 10px; border-radius: 10px; font-weight: 700; }"
        )
        # Set content and color based on test_mode
        if self.test_mode:
            self.mode_badge.setText("TEST MODE")
            self.mode_badge.setStyleSheet(
                "QLabel { background:#b00020; color:white; padding:3px 10px; border-radius:10px; font-weight:700; }"
            )
        else:
            self.mode_badge.setText("STUDY MODE")
            self.mode_badge.setStyleSheet(
                "QLabel { background:#0a7cff; color:white; padding:3px 10px; border-radius:10px; font-weight:700; }"
            )

        self.flag_list_button.clicked.connect(self._open_flag_list)
        head.addWidget(QLabel("Theme: ")); head.addWidget(self.theme_combo)
        head.addStretch(1); head.addWidget(self.timer_label); head.addSpacing(12); 
        head.addWidget(self.progress_label); head.addSpacing(12); head.addWidget(self.mode_badge)
        head.addSpacing(12); head.addWidget(self.flag_button); head.addWidget(self.flag_list_button)
        root.addLayout(head)

        # Question
        self.question_label = QLabel("")
        self.question_label.setWordWrap(True)
        self.question_label.setStyleSheet("font-size:16px; font-weight:600;")
        root.addWidget(self.question_label)

        # Thumbnail under the question (click to open full image)
        self.thumb_label = QLabel()
        self.thumb_label.setVisible(False)  # hidden unless an image exists
        self.thumb_label.setScaledContents(True)
        self.thumb_label.setFixedWidth(60)
        self.thumb_label.setFixedHeight(60)  
        self.thumb_label.setStyleSheet(
            "QLabel { border: 1px solid #e2e2e2; border-radius: 6px; background:#fafafa; }"
        )
        self.thumb_label.setCursor(Qt.CursorShape.PointingHandCursor)

        def _thumb_click(_e):
            self._show_image_for_current()
        self.thumb_label.mousePressEvent = _thumb_click

        root.addWidget(self.thumb_label)

        # Answers
        self.answers_scroll = QScrollArea()
        self.answers_scroll.setWidgetResizable(True)
        self.answers_widget = QWidget()
        self.answers_layout = QVBoxLayout(self.answers_widget)
        self.answers_layout.setContentsMargins(0,0,0,0)
        self.answers_layout.setSpacing(6)
        self.answers_scroll.setWidget(self.answers_widget)
        root.addWidget(self.answers_scroll, 1)

        # Actions row
        actions = QHBoxLayout()
        actions.setSpacing(8)
        self.show_image_button = QPushButton("Show Image")
        self.show_image_button.clicked.connect(self._show_image_for_current)

        # Image badge (clickable pill)
        self.image_badge = QLabel("Image available")
        self.image_badge.setVisible(False)  # hidden until we detect an image
        self.image_badge.setAlignment(Qt.AlignmentFlag.AlignCenter)
        self.image_badge.setStyleSheet(
            """
            QLabel {
                background-color: #0a7cff;
                color: white;
                padding: 3px 8px;
                border-radius: 10px;
                font-size: 11px;
                font-weight: 600;
            }
            QLabel:hover {
                background-color: #086ad6;
            }
            """
        )

        # Make the badge act like a button (open the image dialog)
        def _open_image_from_badge(_e):
            self._show_image_for_current()
        self.image_badge.mousePressEvent = _open_image_from_badge

        self.reason_button = QPushButton("Why?")
        self.reason_button.setToolTip("Show explanation for the answer")
        self.reason_button.clicked.connect(self._show_reason_for_current)
        if self.allow_calc:
            self.calc_button = QPushButton("Calculator")
            self.calc_button.clicked.connect(self._open_calculator)
        self.check_button = QPushButton("Check Answer")
        self.check_button.clicked.connect(self._check_current_answer)
        self.break_button = QPushButton("Take 15-min Break")
        self.break_button.clicked.connect(self._take_break)
        actions.addWidget(self.show_image_button)
        actions.addWidget(self.image_badge)
        actions.addWidget(self.reason_button)
        actions.addWidget(self.check_button)
        if self.allow_calc:
            actions.addWidget(self.calc_button)
        actions.addStretch(1)
        actions.addWidget(self.break_button)
        root.addLayout(actions)

        # Navigation
        nav = QHBoxLayout()
        self.prev_button = QPushButton("Previous"); self.prev_button.clicked.connect(self._prev)
        self.next_button = QPushButton("Next"); self.next_button.clicked.connect(self._next)
        self.submit_button = QPushButton("Submit"); self.submit_button.clicked.connect(self._submit_current)
        self._submit_base_text = self.submit_button.text()
        self._base_button_style = self.submit_button.styleSheet()  # often ""
        self.finish_button = QPushButton("Finish"); self.finish_button.clicked.connect(self._finish)
        nav.addWidget(self.prev_button); nav.addWidget(self.next_button)
        nav.addStretch(1); nav.addWidget(self.submit_button); nav.addWidget(self.finish_button)
        root.addLayout(nav)

        # Start timer if needed
        if self.remaining_seconds > 0:
            self._update_timer_label()
            self.timer.start(1000)

        # Disable features in Test Mode
        if self.test_mode:
            self.check_button.setEnabled(False)
            self.reason_button.setEnabled(False)

        # Break availability
        self._update_break_enabled()

        # Render first question
        self.current_index = 0
        self.score = 0
        self._render_current_question()

    # ---- Helpers ----
    def _update_progress_label(self):
        try:
            total = len(self.quiz)
        except Exception:
            total = 0
        cur = (self.current_index + 1) if total else 0
        self.progress_label.setText(f"{cur}/{total}")

    def _toggle_theme(self):
        app = QApplication.instance()
        current = load_theme_pref()
        new_theme = "dark" if current == "tokyo_night" else "tokyo_night"
        apply_theme(app, new_theme)

        # reapply any per-button override colors (flag/image states)
        self._update_action_buttons_state()
        # if you colorize flag button when flagged, also:
        if hasattr(self, "flags") and self.current_index in self.flags:
            # your existing flagged color override will re-run in _render_current_question
            self._render_current_question()

    def _on_theme_changed(self, index: int):
        key = self.theme_combo.itemData(index)
        if not key:
            return
        app = QApplication.instance()
        apply_theme(app, key)
        # keep your per-button overrides in sync
        self._update_action_buttons_state()
        if hasattr(self, "flags") and self.current_index in self.flags:
            self._render_current_question()

    def _flash_button(self, button: QPushButton, ok: bool = True, ms: int = 900, flash_code = "submit"):
        # remember base text/style once, on the button itself
        if button.property("_base_text") is None:
            button.setProperty("_base_text", button.text())
        if button.property("_base_style") is None:
            button.setProperty("_base_style", button.styleSheet())

        # cancel an in-flight flash (if any)
        t = button.property("_flash_timer")
        if isinstance(t, QTimer):
            t.stop()

        if flash_code == "submit":
            button.setText("✓ Saved")
            button.setStyleSheet("background-color: #9ACD32; color: black;")
        elif flash_code == "check_not_quite":
            button.setText("Not Quite!")
            button.setStyleSheet("background-color: #FFBF00; color: black;")
        elif flash_code == "cc":
            button.setText("Correct!")
            button.setStyleSheet("background-color: #9ACD32; color: black;")

        # start a fresh timer to revert
        timer = QTimer(self)
        timer.setSingleShot(True)

        def _revert():
            base_text = button.property("_base_text") or ""
            base_style = button.property("_base_style") or ""
            button.setText(base_text)
            button.setStyleSheet(base_style)
            button.setProperty("_flash_timer", None)

        timer.timeout.connect(_revert)
        button.setProperty("_flash_timer", timer)
        timer.start(ms)


    def _apply_quiz_settings(self, quiz_data: List[Dict], num_questions: int, allow_repeats: bool) -> List[Dict]:
        if not quiz_data:
            return []
        if num_questions <= 0:
            q = quiz_data[:]; random.shuffle(q); return q
        if num_questions <= len(quiz_data):
            return random.sample(quiz_data, k=num_questions)
        if allow_repeats:
            return [random.choice(quiz_data) for _ in range(num_questions)]
        # cap if no repeats
        return random.sample(quiz_data, k=len(quiz_data))
    

    def _reset_button(self, button: QPushButton):
        t = button.property("_flash_timer")
        if isinstance(t, QTimer):
            t.stop()
            button.setProperty("_flash_timer", None)
        base_text = button.property("_base_text") or button.text()
        base_style = button.property("_base_style") or ""
        button.setText(base_text)
        button.setStyleSheet(base_style)


    # ---- Timer & Breaks ----
    def _update_timer_label(self):
        mm, ss = divmod(max(0, self.remaining_seconds), 60)
        self.timer_label.setText(f"{mm:02d}:{ss:02d}")

    def _on_tick(self):
        self.remaining_seconds -= 1
        if self.remaining_seconds <= 0:
            self.timer.stop()
            self.remaining_seconds = 0
            self._update_timer_label()
            QMessageBox.information(self, "Time", "Time is up! Submitting current answer.")
            self._submit_current()
        else:
            self._update_timer_label()

    def _update_break_enabled(self):
        self.break_button.setEnabled(bool(self.allow_breaks and (self.total_seconds > 0) and (not self.break_taken)))

    def _take_break(self):
        if self.break_taken or self.total_seconds <= 0 or not self.allow_breaks:
            return

        resp = QMessageBox.question(
            self,
            "Take Break",
            "Start a single 15-minute break? The test timer will be paused.",
            QMessageBox.StandardButton.Yes | QMessageBox.StandardButton.No
        )
        if resp != QMessageBox.StandardButton.Yes:
            return

        # Pause main test timer
        was_running = self.timer.isActive()
        if was_running:
            self.timer.stop()

        # Show countdown dialog
        dialog = BreakDialog(self, total_seconds=15 * 60)
        dialog.exec()

        # Mark the one allowed break as used
        self.break_taken = True
        self._update_break_enabled()

        # Resume main timer if there’s time left
        if was_running and self.remaining_seconds > 0:
            self.timer.start(1000)


    def _set_thumbnail_for_current(self):
        """Show/hide the thumbnail based on the current question's image."""
        if not self.quiz:
            self.thumb_label.setVisible(False)
            return
        image = self.quiz[self.current_index].get("image")
        if image and os.path.exists(image):
            pix = QPixmap(image)
            if not pix.isNull():
                # scale keeping aspect ratio to the fixed height we set on the label
                scaled = pix.scaledToHeight(self.thumb_label.height(), Qt.TransformationMode.SmoothTransformation)
                self.thumb_label.setPixmap(scaled)
                self.thumb_label.setToolTip("Click to view full image")
                self.thumb_label.setVisible(True)
                return
        # no/invalid image → hide thumbnail
        self.thumb_label.clear()
        self.thumb_label.setVisible(False)
    
    # ---- Navigation ----
    def _prev(self):
        self._save_current_selection()
        if self.current_index > 0:
            self.current_index -= 1
            self._render_current_question()
            self._update_progress_label()

    def _next(self):
        self._save_current_selection()
        if self.current_index < len(self.quiz) - 1:
            self.current_index += 1
            self._render_current_question()
            self._update_progress_label()

    # ---- Flagging ----

    def _toggle_flag(self):
        idx = self.current_index
        if idx in self.flags:
            self.flags.remove(idx)
            self.flag_button.setStyleSheet(
                "QPushButton {"
                " background-color: #b00020; color: white;"
                " padding: 6px 12px; border-radius: 6px;"
                " font-weight: 600;"
                "}"
            )
        else:
            self.flags.add(idx)
            self.flag_button.setStyleSheet("")
        self._render_current_question()  # to update header indicator

    def _open_flag_list(self):
        if not self.flags:
            QMessageBox.information(self, "Flagged", "No flagged questions.")
            return
        dialog = QDialog(self); dialog.setWindowTitle("Flagged Questions")
        v = QVBoxLayout(dialog)
        lst = QListWidget()
        for i in sorted(self.flags):
            text = f"{i+1}. {self.quiz[i].get('question','')[:80]}"
            item = QListWidgetItem(text)
            item.setData(Qt.ItemDataRole.UserRole, i)
            lst.addItem(item)
        v.addWidget(lst)
        row = QHBoxLayout()
        go = QPushButton("Go To"); close = QPushButton("Close")
        row.addStretch(1); row.addWidget(go); row.addWidget(close)
        v.addLayout(row)
        def _go():
            it = lst.currentItem()
            if not it: return
            self.current_index = it.data(Qt.ItemDataRole.UserRole)
            dialog.accept(); self._render_current_question()
        go.clicked.connect(_go); close.clicked.connect(dialog.reject)
        dialog.exec()

    # ---- Rendering / selection ----
    def _make_radio_row(self, option_text: str, group: QButtonGroup) -> QWidget:
        row = QWidget(); row.setSizePolicy(QSizePolicy.Policy.Expanding, QSizePolicy.Policy.Maximum)
        hl = QHBoxLayout(row); hl.setContentsMargins(0,0,0,0); hl.setSpacing(6)
        button = QRadioButton(); button.setProperty("optionText", option_text); button.setToolTip(option_text)
        group.addButton(button)
        label = QLabel(option_text); label.setWordWrap(True)
        label.setTextInteractionFlags(Qt.TextInteractionFlag.TextSelectableByMouse)
        label.setSizePolicy(QSizePolicy.Policy.Expanding, QSizePolicy.Policy.Preferred)
        def click_label(_e): button.click()
        label.mousePressEvent = click_label
        hl.addWidget(button, 0, Qt.AlignmentFlag.AlignTop); hl.addWidget(label, 1)
        return row

    def _make_check_row(self, option_text: str) -> Tuple[QCheckBox, QWidget]:
        row = QWidget(); row.setSizePolicy(QSizePolicy.Policy.Expanding, QSizePolicy.Policy.Maximum)
        hl = QHBoxLayout(row); hl.setContentsMargins(0,0,0,0); hl.setSpacing(6)
        cb = QCheckBox(); cb.setProperty("optionText", option_text); cb.setToolTip(option_text)
        label = QLabel(option_text); label.setWordWrap(True)
        label.setTextInteractionFlags(Qt.TextInteractionFlag.TextSelectableByMouse)
        label.setSizePolicy(QSizePolicy.Policy.Expanding, QSizePolicy.Policy.Preferred)
        def click_label(_e): cb.toggle()
        label.mousePressEvent = click_label
        hl.addWidget(cb, 0, Qt.AlignmentFlag.AlignTop); hl.addWidget(label, 1)
        return cb, row

    def _render_current_question(self):
        # Clear answers layout
        while self.answers_layout.count():
            item = self.answers_layout.takeAt(0)
            w = item.widget()
            if w: w.deleteLater()

        if not self.quiz:
            self.question_label.setText("No questions loaded.")
            return

        q = self.quiz[self.current_index]
        text = q.get("question","")
        multi = bool(q.get("multi", False))
        has_image = bool(self.quiz[self.current_index].get("image"))
        self.question_label.setText(f"{self.current_index+1}. {text}")
        # Update thumbnail (if any) for this question
        self._set_thumbnail_for_current()

        opts = q.get("options", []) or []
        if multi:
            hint = QLabel("(Select all that apply)")
            hint.setObjectName("MultiHint")
            hint.setWordWrap(True)
            self.answer_group = None
            self.checkboxes: List[QCheckBox] = []
            self.answers_layout.addWidget(hint, 0, Qt.AlignmentFlag.AlignTop)
            for opt in opts:
                cb, row = self._make_check_row(opt)
                self.checkboxes.append(cb)
                self.answers_layout.addWidget(row, 0, Qt.AlignmentFlag.AlignTop)
            prev = self.user_answers[self.current_index]
            if isinstance(prev, set):
                for cb in self.checkboxes:
                    cb.setChecked(cb.property("optionText") in prev if cb.property("optionText") else False)
        else:
            self.checkboxes = []
            self.answer_group = QButtonGroup(self)
            for opt in opts:
                row = self._make_radio_row(opt, self.answer_group)
                self.answers_layout.addWidget(row, 0, Qt.AlignmentFlag.AlignTop)
            prev = self.user_answers[self.current_index]
            if isinstance(prev, str):
                for b in self.answer_group.buttons():
                    if b.property("optionText") == prev:
                        b.setChecked(True); break
                    
        # Keep flag button style in sync with current question
        if self.current_index in self.flags:
            self.flag_button.setStyleSheet(
                "QPushButton { background-color: #b00020; color: white; padding: 6px 12px; border-radius: 6px; font-weight: 600; }"
                "QPushButton:hover { background-color: #d32f2f; }"
            )
        else:
            self.flag_button.setStyleSheet("")
        # Ensure Submit button is back to normal on each question render
        self._update_progress_label()
        self._reset_button(self.submit_button)
        self._reset_button(self.check_button)
        self.answers_scroll.verticalScrollBar().setValue(0)
        self._update_action_buttons_state()

    def _update_action_buttons_state(self):
        has_image = False
        if self.quiz:
            image = self.quiz[self.current_index].get("image")
            has_image = bool(image and os.path.exists(image))
        # Enable/disable and colorize the Show Image button
        self.show_image_button.setEnabled(has_image)
        self.show_image_button.setText("Show Image")
        if has_image:
            # Primary color when available
            self.show_image_button.setStyleSheet(
                "QPushButton {"
                " background-color: #0a7cff; color: white;"
                " padding: 6px 12px; border-radius: 6px;"
                "}"
            )
        else:
            # Subtle/disabled look when no image
            self.show_image_button.setStyleSheet("")


    def _save_current_selection(self):
        if not self.quiz:
            return
        q = self.quiz[self.current_index]
        multi = bool(q.get("multi", False))
        if multi and hasattr(self, "checkboxes"):
            chosen = set()
            for cb in self.checkboxes:
                if cb.isChecked():
                    text = cb.property("optionText")
                    if text: chosen.add(text)
            self.user_answers[self.current_index] = chosen
        elif hasattr(self, "answer_group") and self.answer_group:
            button = self.answer_group.checkedButton()
            self.user_answers[self.current_index] = button.property("optionText") if button else None

    # ---- Actions ----
    def _submit_current(self):
        if not self.quiz:
            return
        self._save_current_selection()
        # Neutral behavior in Test Mode; no correctness reveal
        self._flash_button(self.submit_button, ok=True)

        if self.test_mode:
            return

        # Normal mode: show correctness
        q = self.quiz[self.current_index]
        correct = q.get("answer", set())
        if isinstance(correct, list):
            correct = set(correct)
        elif isinstance(correct, str):
            correct = {correct}
        chosen = self.user_answers[self.current_index]
        is_correct = False
        if isinstance(chosen, set):
            is_correct = chosen == set(correct)
        elif isinstance(chosen, str):
            is_correct = chosen in correct

        self._flash_button(self.submit_button, ok=is_correct, flash_code="submit")


    def _check_current_answer(self):
        if self.test_mode:
            return
        if not self.quiz: return
        self._save_current_selection()
        q = self.quiz[self.current_index]
        correct = q.get("answer", set())
        if isinstance(correct, list): correct = set(correct)
        elif isinstance(correct, str): correct = {correct}
        chosen = self.user_answers[self.current_index]
        is_correct = False
        if isinstance(chosen, set):
            is_correct = chosen == set(correct)
        elif isinstance(chosen, str):
            is_correct = chosen in correct
        if is_correct:
            self._flash_button(self.check_button, ok=True, ms=900, flash_code="check_current")
        else:
            self._flash_button(self.check_button, ok=True, ms=900, flash_code="check_not_quite")


    def _show_image_for_current(self):
        q = self.quiz[self.current_index]
        image = q.get("image")
        if image and os.path.exists(image):
            dialog = ImageViewerDialog([image], parent=self, title="Question Image")
            dialog.exec()
        else:
            QMessageBox.information(self, "Image", "No image for this question.")

    def _show_reason_for_current(self):
        if self.test_mode:
            return
        q = self.quiz[self.current_index]
        reason = q.get("explanation", "") or "No explanation provided."
        image = q.get("image")
        # simple dialog
        dialog = QDialog(self); dialog.setWindowTitle("Explanation")
        v = QVBoxLayout(dialog)
        if image and os.path.exists(image):
            pix = QPixmap(image); label_image = QLabel(); label_image.setPixmap(pix) 
            label_image.setScaledContents(True); label_image.setMinimumHeight(min(240, pix.height()))
            v.addWidget(label_image)
        label = QLabel(reason); label.setWordWrap(True); v.addWidget(label)
        row = QHBoxLayout(); row.addStretch(1); ok = QPushButton("Close"); ok.clicked.connect(dialog.accept); row.addWidget(ok)
        v.addLayout(row)
        dialog.resize(640, 480); dialog.exec()

    def _open_calculator(self):
        CalculatorPopup(self).exec()

    def _open_review_dialog(self):
        items = []
        correct_count = 0
        total = len(self.quiz)

        for i, q in enumerate(self.quiz, start=1):
            # normalize the correct answers into a set
            correct = q.get("answer", set())
            if isinstance(correct, list):
                correct = set(correct)
            elif isinstance(correct, str):
                correct = {correct}
            elif not isinstance(correct, set):
                correct = set()

            chosen_raw = self.user_answers[i - 1]
            if isinstance(chosen_raw, set):
                is_correct_now = (chosen_raw == correct)
                chosen_str = ", ".join(sorted(chosen_raw))
            elif chosen_raw is None:
                is_correct_now = False
                chosen_str = "(no answer)"
            else:
                is_correct_now = (chosen_raw in correct)
                chosen_str = str(chosen_raw)

            if is_correct_now:
                correct_count += 1

            correct_str = ", ".join(sorted(correct))
            items.append({
                "question": q.get("question", ""),
                "correct": correct_str,
                "chosen": chosen_str,
                "explanation": q.get("explanation", ""),
                "flagged": (i - 1) in self.flags
            })

        percent = int(round((correct_count / total) * 100)) if total else 0
        score_tuple = (correct_count, total, percent)

        # Restart logic for new run
        def _restart():
            data = self.raw_quiz
            settings = start_with_settings_dialog(None, data, pptx_path=self.pptx_path, return_settings=True)
            if settings is None:
                return
            new_win = QuizMainWindow(data, pptx_path=self.pptx_path, settings=settings)
            new_win.show()

        dialog = ReviewPopup(
            None,
            items,
            score_tuple,
            pptx_basename=(self.pptx_path or "review"),
            restart_callback=_restart
        )
        dialog.exec()


    def _finish(self):
        # Always clickable; warn if unanswered
        self._save_current_selection()
        unanswered = any(ans is None or (isinstance(ans, set) and not ans) or (isinstance(ans, str) and not ans.strip())
                         for ans in self.user_answers)
        if unanswered:
            resp = QMessageBox.question(self, "Finish quiz?",
                                        "Some questions are unanswered. Finish anyway and see the review?",
                                        QMessageBox.StandardButton.Yes | QMessageBox.StandardButton.No,
                                        QMessageBox.StandardButton.No)
            if resp != QMessageBox.StandardButton.Yes:
                return
        # Close main and open review
        self.close()
        self._open_review_dialog()

# -----------------------------
# App flow
# -----------------------------

def apply_quiz_settings(quiz_data: List[Dict],
                        num_questions: int = 0,
                        allow_repeats: bool = True) -> List[Dict]:
    if not quiz_data:
        return []
    if num_questions <= 0 or num_questions >= len(quiz_data):
        q = quiz_data[:]
        random.shuffle(q)
        return q
    if allow_repeats and num_questions > len(quiz_data):
        return [random.choice(quiz_data) for _ in range(num_questions)]
    return random.sample(quiz_data, k=num_questions)


def start_with_settings_dialog(parent: Optional[QWidget],
                               quiz_data: List[Dict],
                               pptx_path: Optional[str] = None,
                               return_settings: bool = False):
    if not quiz_data:
        QMessageBox.warning(parent, "No Questions", "No questions were loaded.")
        return None
    dialog = QuestionPopup(parent, max_questions=len(quiz_data))
    if dialog.exec():
        result = dialog.get_result()  # (count, timer_seconds, allow_repeats, allow_calc, test_mode, allow_breaks)
        if not result:
            return None
        count, timer_seconds, allow_repeats, allow_calc, test_mode, allow_breaks = result
        if return_settings:
            return (count, timer_seconds, allow_repeats, allow_calc, test_mode, allow_breaks)
        window = QuizMainWindow(quiz_data, pptx_path=pptx_path, settings=result)
        window.show()
        return window
    return None


def main_open_pptx_and_run():
    app = QApplication(sys.argv)
    # Base theme for all QPushButtons and labels
    apply_theme(app, load_theme_pref())

    # Ask for PPTX
    dialog = QFileDialog()
    dialog.setFileMode(QFileDialog.FileMode.ExistingFile)
    dialog.setNameFilter("PowerPoint files (*.pptx);;All files (*.*)")
    if not dialog.exec():
        return
    selected = dialog.selectedFiles()
    if not selected:
        return
    pptx_file = selected[0]
    # Build quiz
    quiz_data, pptx_basename = build_quiz_from_pptx(pptx_file)
    if not quiz_data:
        QMessageBox.warning(None, "No Questions", "Could not find any questions in that PPTX.")
        return
    # Start with settings
    window = start_with_settings_dialog(None, quiz_data, pptx_path=pptx_basename)
    if not window:
        return
    sys.exit(app.exec())


def run_quiz_app(quiz_data: List[Dict], pptx_path: Optional[str] = None, *, timer_seconds: int = 0):
    app = QApplication(sys.argv)
    settings = (0, timer_seconds, True, False, True)  # defaults
    window = QuizMainWindow(quiz_data, pptx_path=pptx_path, settings=settings)
    if timer_seconds and timer_seconds > 0:
        window.timer.start(1000)
    window.show()
    sys.exit(app.exec())


if __name__ == "__main__":
    ensure_requirements()
    main_open_pptx_and_run()
